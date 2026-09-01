"""Assignment Auto-Worker — AI mengerjakan tugas e-learning & menghasilkan dokumen siap kumpul.

Alur:
1. Ambil detail tugas dari e-learning (deskripsi/instruksi + lampiran soal).
2. Ekstrak teks lampiran + materi kuliah terkait sebagai konteks.
3. LLM mengerjakan tugas mengikuti format/instruksi yang diminta soal.
4. Render hasil ke file .docx rapi (header nama/NIM, heading, bullet, kode program).
5. Catat hasil ke data/metadata/assignment_outputs.json — file-nya dikirim ke
   Telegram untuk direview; PENGUMPULAN TETAP MANUAL oleh mahasiswa.
"""
import json
import logging
import re
import zipfile
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Pt

from src.ai_service import AIService
from src.config import (
    ASSIGNMENTS_ATTACH_DIR,
    ASSIGNMENTS_OUTPUT_DIR,
    ASSIGNMENTS_FILE,
    ASSIGNMENT_OUTPUTS_FILE,
    EXTRACTED_TEXT_DIR,
    STUDENT_NAME,
    STUDENT_NIM,
    UMN_USERNAME
)
from src.moodle_client import MoodleClient

logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")
logger = logging.getLogger("assignment_worker")

MAX_ATTACHMENT_CHARS = 20000   # total teks lampiran
MAX_PER_FILE_CHARS = 12000     # teks per lampiran
MAX_CONTEXT_CHARS = 20000      # materi kuliah


class AssignmentWorker:
    def __init__(self, client: Optional[MoodleClient] = None):
        self.ai = AIService()
        self._client = client

    # ---------------------------------------------------------------- utils
    def _get_client(self) -> MoodleClient:
        if self._client is None:
            self._client = MoodleClient()
        if not self._client.is_logged_in:
            self._client.login()
        return self._client

    def _student_name(self) -> str:
        if STUDENT_NAME:
            return STUDENT_NAME
        # Turunkan dari email: "kenny.valent@student.umn.ac.id" -> "Kenny Valent"
        local = UMN_USERNAME.split("@")[0] if "@" in UMN_USERNAME else UMN_USERNAME
        return local.replace(".", " ").replace("_", " ").title() or "Mahasiswa UMN"

    def _student_nim(self) -> str:
        return STUDENT_NIM or "-"

    @staticmethod
    def _safe_name(name: str) -> str:
        clean = re.sub(r'[\\/*?:"<>|]', "_", name).strip()
        return re.sub(r"\s+", " ", clean)[:90]

    # ---------------------------------------------------------------- state
    def list_pending(self) -> List[Dict[str, Any]]:
        """Tugas pending (belum disubmit) dari assignments.json."""
        if not ASSIGNMENTS_FILE.exists():
            return []
        try:
            assignments = json.loads(ASSIGNMENTS_FILE.read_text(encoding="utf-8"))
            return [a for a in assignments if a and not a.get("is_submitted")]
        except Exception as e:
            logger.warning(f"Gagal baca assignments.json: {e}")
            return []

    def list_done_urls(self) -> List[str]:
        if not ASSIGNMENT_OUTPUTS_FILE.exists():
            return []
        try:
            data = json.loads(ASSIGNMENT_OUTPUTS_FILE.read_text(encoding="utf-8"))
            return list(data.keys())
        except Exception:
            return []

    def _record_output(self, url: str, record: Dict[str, Any]) -> None:
        try:
            data = {}
            if ASSIGNMENT_OUTPUTS_FILE.exists():
                data = json.loads(ASSIGNMENT_OUTPUTS_FILE.read_text(encoding="utf-8"))
            data[url] = record
            ASSIGNMENT_OUTPUTS_FILE.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")
        except Exception as e:
            logger.warning(f"Gagal mencatat output tugas: {e}")

    # ---------------------------------------------------------- ekstraksi teks
    @staticmethod
    def _extract_file_text(path: Path) -> str:
        """Ekstrak teks dari lampiran soal: PDF, DOCX, PPTX, TXT."""
        try:
            suffix = path.suffix.lower()
            if suffix == ".pdf":
                from pypdf import PdfReader
                reader = PdfReader(str(path))
                text = "\n".join((page.extract_text() or "") for page in reader.pages)
            elif suffix == ".docx":
                from docx import Document as DocxDocument
                doc = DocxDocument(str(path))
                parts = [p.text for p in doc.paragraphs]
                for table in doc.tables:
                    for row in table.rows:
                        parts.append(" | ".join(c.text.strip() for c in row.cells))
                text = "\n".join(parts)
            elif suffix == ".pptx":
                from pptx import Presentation
                prs = Presentation(str(path))
                parts = []
                for i, slide in enumerate(prs.slides, 1):
                    parts.append(f"[Slide {i}]")
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                t = "".join(run.text for run in para.runs).strip()
                                if t:
                                    parts.append(t)
                text = "\n".join(parts)
            elif suffix in (".txt", ".md", ".csv"):
                text = path.read_text(encoding="utf-8", errors="ignore")
            else:
                return ""
            return text[:MAX_PER_FILE_CHARS]
        except Exception as e:
            logger.warning(f"Gagal ekstrak teks {path.name}: {e}")
            return ""

    def _course_context(self, course_name: str, max_chars: int = MAX_CONTEXT_CHARS) -> str:
        """Kumpulkan teks materi kuliah terkait (folder extracted_text bernama sama dengan kursus)."""
        folder = EXTRACTED_TEXT_DIR / course_name
        if not folder.is_dir():
            # fallback: cari folder yang mengandung kode kursus, mis. (IF571-B)
            m = re.search(r"\(([A-Z]{2,6}\s?\d{3,4})", course_name)
            if m:
                code = re.sub(r"[^A-Z0-9]", "", m.group(1)).lower()
                folder = next((d for d in EXTRACTED_TEXT_DIR.iterdir()
                               if d.is_dir() and code in d.name.lower()), None)
        if folder is None or not folder.is_dir():
            return "(Materi kuliah untuk mata kuliah ini belum tersedia.)"

        blocks, total = [], 0
        for f in sorted(folder.glob("*.txt")):
            try:
                content = f.read_text(encoding="utf-8", errors="ignore")
            except Exception:
                continue
            piece = f"\n--- MATERI: {f.name} ---\n{content}\n"
            if total + len(piece) > max_chars:
                remaining = max_chars - total
                if remaining > 500:
                    blocks.append(piece[:remaining] + "\n[TRUNCATED]")
                break
            blocks.append(piece)
            total += len(piece)
        return "\n".join(blocks) or "(Materi kuliah masih kosong.)"

    # ------------------------------------------------------------------ AI
    def _build_prompt(self, assignment: Dict[str, Any], description: str,
                      attachment_text: str, context: str) -> str:
        name, nim = self._student_name(), self._student_nim()
        return f"""Kamu adalah mahasiswa UMN bernama {name} (NIM {nim}) yang rajin dan pandai.
Kerjakan tugas kuliah berikut SEBAIK MUNGKIN dan SESUAI FORMAT YANG DIMINTA SOAL.

=== IDENTITAS TUGAS ===
Mata Kuliah : {assignment.get('course_name', '-')}
Judul Tugas : {assignment.get('title', '-')}
Deadline    : {assignment.get('due_date', '-')}

=== INSTRUKSI / SOAL TUGAS (dari e-learning) ===
{description or '(Deskripsi kosong — kerjakan berdasarkan judul tugas dan materi kuliah yang relevan.)'}

=== LAMPIRAN SOAL (teks hasil ekstraksi) ===
{attachment_text or '(Tidak ada lampiran soal.)'}

=== MATERI KULIAH RELEVAN ===
{context}

=== ATURAN PENGERJAAN ===
1. Ikuti SEMUA instruksi & format dari soal: struktur bab/bagian, jumlah kata/kalimat, bahasa
   (kalau soal tidak menyebutkan bahasa, pakai Bahasa Indonesia akademik yang rapi).
2. Jawaban harus spesifik ke materi kuliah ini — kutip konsep/istilah dari materi, jangan
   jawaban generik template.
3. JANGAN mengarang data statistik, kutipan buku, atau sumber referensi palsu. Kalau butuh
   referensi, gunakan materi kuliah yang diberikan.
4. Kalau soal meminta kode program, tulis kode yang benar dan lengkap di field "code" bagian
   terkait, plus penjelasan singkat di "paragraphs".
5. Keluarkan HANYA JSON valid (tanpa teks lain, tanpa ```json fence) dengan struktur persis:
{{
  "summary": "ringkasan 3-5 kalimat: apa yang dikerjakan dan bagaimana strukturnya",
  "files": [
    {{
      "filename": "Tugas - <judul singkat>.docx",
      "sections": [
        {{
          "heading": "Judul Bagian 1",
          "paragraphs": ["paragraf isi..."],
          "bullets": ["poin 1", "poin 2"],
          "code": "kode program bila ada, string biasa dengan \\n untuk baris baru, atau string kosong"
        }}
      ]
    }}
  ]
}}
Ketentuan: cukup 1 file DOCX. Setiap "sections" punya minimal satu dari paragraphs/bullets/code.
Jangan pakai markdown (**, ##) di dalam teks — sudah diformat otomatis jadi dokumen Word."""

    @staticmethod
    def _parse_llm_json(text: str) -> Optional[Dict[str, Any]]:
        """Parse JSON dari balasan LLM dengan toleransi: code fence, teks di luar JSON,
        dan newline/control character literal di dalam string (kesalahan umum LLM)."""
        if not text:
            return None
        cleaned = re.sub(r"```(?:json)?", "", text).strip().strip("`")
        candidates = [cleaned]
        start, end = cleaned.find("{"), cleaned.rfind("}")
        if start != -1 and end > start:
            candidates.append(cleaned[start:end + 1])
        for cand in candidates:
            for strict in (True, False):  # strict=False: izinkan \n literal di string
                try:
                    parsed = json.loads(cand, strict=strict)
                    if isinstance(parsed, dict):
                        return parsed
                except Exception:
                    continue
        return None

    # ------------------------------------------------------------ render docx
    def _render_docx(self, spec: Dict[str, Any], assignment: Dict[str, Any], out_path: Path) -> Path:
        doc = Document()

        # Sampul / header identitas
        title_para = doc.add_paragraph()
        title_run = title_para.add_run(spec.get("filename", assignment.get("title", "Tugas")).replace(".docx", ""))
        title_run.bold = True
        title_run.font.size = Pt(16)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

        meta_para = doc.add_paragraph()
        meta = (f"{assignment.get('course_name', '')}\n"
                f"Nama : {self._student_name()}\n"
                f"NIM  : {self._student_nim()}\n"
                f"Tanggal : {datetime.now():%d %B %Y}")
        meta_run = meta_para.add_run(meta)
        meta_run.font.size = Pt(10)
        meta_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph()

        for section in spec.get("sections", []):
            heading = (section.get("heading") or "").strip()
            if heading:
                doc.add_heading(heading, level=1)
            for para in section.get("paragraphs", []) or []:
                p = doc.add_paragraph(str(para))
                p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            for bullet in section.get("bullets", []) or []:
                doc.add_paragraph(str(bullet), style="List Bullet")
            code = (section.get("code") or "").strip()
            if code:
                code_heading = doc.add_paragraph()
                r = code_heading.add_run("Kode Program:")
                r.bold = True
                for line in code.split("\n"):
                    cp = doc.add_paragraph()
                    run = cp.add_run(line)
                    run.font.name = "Courier New"
                    run.font.size = Pt(9)
                    cp.paragraph_format.space_after = Pt(0)

        doc.save(str(out_path))
        return out_path

    # ------------------------------------------------------------- main entry
    def work_on_assignment(self, assignment: Dict[str, Any]) -> Dict[str, Any]:
        """Kerjakan satu tugas: ambil soal → AI → DOCX → kembalikan hasil + file."""
        url = assignment.get("url", "")
        title = assignment.get("title", "Tugas")
        course_name = assignment.get("course_name", "")
        result: Dict[str, Any] = {"ok": False, "assignment": title, "files": [], "summary": "", "error": ""}

        try:
            client = self._get_client()
            logger.info(f"Mengerjakan tugas: {title} ({course_name})")
            details = client.get_assignment_details(url, course_name, title)

            # Teks lampiran soal
            attach_parts = []
            total_attach = 0
            for att in details.get("attachments", []):
                text = self._extract_file_text(Path(att["path"]))
                if not text:
                    continue
                piece = f"\n--- LAMPIRAN: {Path(att['path']).name} ---\n{text}\n"
                if total_attach + len(piece) > MAX_ATTACHMENT_CHARS:
                    remaining = MAX_ATTACHMENT_CHARS - total_attach
                    if remaining > 500:
                        attach_parts.append(piece[:remaining] + "\n[TRUNCATED]")
                    break
                attach_parts.append(piece)
                total_attach += len(piece)
            attachment_text = "\n".join(attach_parts)

            context = self._course_context(course_name)
            prompt = self._build_prompt(assignment, details.get("description", ""), attachment_text, context)

            llm_text = self.ai._generate_with_fallback(prompt)
            spec = self._parse_llm_json(llm_text)
            if not spec or not spec.get("files"):
                # Fallback: pakai teks mentah sebagai satu dokumen
                logger.warning(f"LLM tidak mengembalikan JSON valid — fallback ke teks mentah. Raw: {llm_text[:200]!r}")
                paragraphs = [l.strip() for l in llm_text.split("\n") if l.strip()]
                spec = {"summary": paragraphs[0][:300] if paragraphs else "",
                        "files": [{"filename": f"{self._safe_name(title)}.docx",
                                   "sections": [{"heading": "Jawaban", "paragraphs": paragraphs}]}]}

            # Render file(s)
            stamp = datetime.now().strftime("%Y%m%d_%H%M")
            out_files: List[Path] = []
            tmp_dir = ASSIGNMENTS_OUTPUT_DIR / f"{stamp}_{self._safe_name(title)[:60]}"
            tmp_dir.mkdir(parents=True, exist_ok=True)

            for f_spec in spec.get("files", []):
                fname = self._safe_name(f_spec.get("filename") or f"{title}.docx")
                if not fname.lower().endswith(".docx"):
                    fname += ".docx"
                out_path = tmp_dir / fname
                self._render_docx(f_spec, assignment, out_path)
                out_files.append(out_path)
                logger.info(f"DOCX dibuat: {out_path}")

            # Kalau ada lebih dari satu file → zip jadi satu paket
            if len(out_files) > 1:
                zip_path = tmp_dir / f"{self._safe_name(title)[:80]}.zip"
                with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
                    for fp in out_files:
                        zf.write(fp, fp.name)
                out_files = [zip_path]

            result.update({
                "ok": True,
                "files": [str(p) for p in out_files],
                "summary": spec.get("summary", ""),
                "description_found": bool(details.get("description")),
                "attachments": [Path(a["path"]).name for a in details.get("attachments", [])]
            })

            self._record_output(url, {
                "title": title,
                "course_name": course_name,
                "generated_at": datetime.now().isoformat(timespec="seconds"),
                "files": result["files"],
                "summary": result["summary"],
                "submitted": False  # pengumpulan tetap manual oleh mahasiswa
            })

        except Exception as e:
            logger.error(f"Gagal mengerjakan tugas '{title}': {e}")
            result["error"] = str(e)

        return result

    def work_on_many(self, assignments: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        return [self.work_on_assignment(a) for a in assignments]
