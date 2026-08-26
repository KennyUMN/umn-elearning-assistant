import os
import shutil
import logging
from pathlib import Path
from typing import Optional, List, Dict

import pypdf
from pptx import Presentation
import docx
from PIL import Image
try:
    import pytesseract
    HAS_TESSERACT = shutil.which("tesseract") is not None
except ImportError:
    HAS_TESSERACT = False

from src.config import MATERIALS_DIR, EXTRACTED_TEXT_DIR

logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")
logger = logging.getLogger("document_parser")

class DocumentParser:
    def __init__(self, materials_dir: Path = MATERIALS_DIR, output_dir: Path = EXTRACTED_TEXT_DIR):
        self.materials_dir = materials_dir
        self.output_dir = output_dir

    def extract_pdf(self, file_path: Path) -> str:
        """Extract text from PDF file page-by-page."""
        text_chunks = []
        try:
            reader = pypdf.PdfReader(str(file_path))
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text() or ""
                page_text = page_text.strip()
                if page_text:
                    text_chunks.append(f"--- [Page {i + 1}] ---\n{page_text}")
        except Exception as e:
            logger.error(f"Error reading PDF {file_path.name}: {e}")
        return "\n\n".join(text_chunks)

    def extract_pptx(self, file_path: Path) -> str:
        """Extract text from PPTX slides, text frames, shapes, and presenter notes."""
        text_chunks = []
        try:
            prs = Presentation(str(file_path))
            for i, slide in enumerate(prs.slides):
                slide_lines = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            line = paragraph.text.strip()
                            if line:
                                slide_lines.append(line)

                # Extract notes if any
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    note = slide.notes_slide.notes_text_frame.text.strip()
                    if note:
                        slide_lines.append(f"[Notes: {note}]")

                if slide_lines:
                    text_chunks.append(f"--- [Slide {i + 1}] ---\n" + "\n".join(slide_lines))
        except Exception as e:
            logger.error(f"Error reading PPTX {file_path.name}: {e}")
        return "\n\n".join(text_chunks)

    def extract_docx(self, file_path: Path) -> str:
        """Extract text from DOCX paragraphs and tables."""
        lines = []
        try:
            doc = docx.Document(str(file_path))
            for p in doc.paragraphs:
                text = p.text.strip()
                if text:
                    lines.append(text)
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_text:
                        lines.append(" | ".join(row_text))
        except Exception as e:
            logger.error(f"Error reading DOCX {file_path.name}: {e}")
        return "\n".join(lines)

    def parse_file(self, file_path: Path) -> Optional[str]:
        """Detect file format and extract textual content."""
        ext = file_path.suffix.lower()
        if ext == ".pdf":
            return self.extract_pdf(file_path)
        elif ext in [".pptx", ".ppt"]:
            return self.extract_pptx(file_path)
        elif ext in [".docx", ".doc"]:
            return self.extract_docx(file_path)
        elif ext in [".txt", ".md"]:
            try:
                return file_path.read_text(encoding="utf-8", errors="ignore")
            except Exception:
                return None
        return None

    def process_all(self) -> List[Dict[str, str]]:
        """Scan all materials directory and convert all files to extracted text."""
        extracted_docs = []
        for course_dir in self.materials_dir.iterdir():
            if not course_dir.is_dir():
                continue

            course_name = course_dir.name
            target_course_text_dir = self.output_dir / course_name
            target_course_text_dir.mkdir(parents=True, exist_ok=True)

            for root, _, files in os.walk(course_dir):
                for fname in files:
                    if fname.startswith("."):
                        continue
                    file_path = Path(root) / fname
                    target_txt_file = target_course_text_dir / f"{file_path.stem}.txt"

                    # Skip if already extracted and source hasn't changed
                    if target_txt_file.exists() and target_txt_file.stat().st_mtime >= file_path.stat().st_mtime:
                        continue

                    logger.info(f"Extracting text from: [{course_name}] {fname}")
                    content = self.parse_file(file_path)
                    if content and content.strip():
                        header = f"=== MATA KULIAH: {course_name} ===\n=== DOKUMEN: {fname} ===\n\n"
                        target_txt_file.write_text(header + content, encoding="utf-8")
                        extracted_docs.append({
                            "course": course_name,
                            "file": fname,
                            "txt_path": str(target_txt_file)
                        })

        logger.info(f"Extraction pass finished. Processed {len(extracted_docs)} new/updated files.")
        return extracted_docs
